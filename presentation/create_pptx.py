"""Generate Segger v0.2.0 Technical Presentation as PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colour palette ──────────────────────────────────────────────
DEEP_BLUE = RGBColor(0, 60, 113)
TEAL_DARK = RGBColor(0, 100, 120)
TEAL_ACCENT = RGBColor(0, 150, 136)
CORAL = RGBColor(231, 76, 60)
GOLD = RGBColor(241, 196, 15)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_TEXT = RGBColor(30, 30, 30)
MED_GRAY = RGBColor(120, 120, 120)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ─────────────────────────────────────────────────────
def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, font_name="Calibri", bold_first=False,
                   bullet=False, line_spacing=1.3):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "\u2022 " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_equation_box(slide, left, top, width, height, text, font_size=16):
    """Add a shaded box for equations (using Cambria Math monospace look)."""
    rect = add_rect(slide, left, top, width, height, RGBColor(235, 245, 255))
    rect.line.color.rgb = RGBColor(180, 210, 240)
    rect.line.width = Pt(1)
    txBox = add_text_box(slide, left + Inches(0.15), top + Inches(0.08),
                         width - Inches(0.3), height - Inches(0.16),
                         text, font_size=font_size, font_name="Cambria Math",
                         color=DEEP_BLUE)
    return txBox


def title_bar(slide, title_text):
    """Add deep blue title bar at top."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), DEEP_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.12), Inches(11), Inches(0.7),
                 title_text, font_size=28, bold=True, color=WHITE)


def section_slide(slide, section_text, subtitle=""):
    """Full-slide section divider."""
    add_bg(slide, DEEP_BLUE)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 section_text, font_size=42, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
                     subtitle, font_size=22, color=TEAL_ACCENT,
                     alignment=PP_ALIGN.CENTER)


def slide_number(slide, num, total=30):
    add_text_box(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}", font_size=11, color=MED_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, DEEP_BLUE)
add_rect(sl, Inches(0), Inches(5.8), W, Inches(0.06), TEAL_ACCENT)

add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Segger v0.2.0", font_size=52, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.0), Inches(11), Inches(1.2),
             "Multi-Task Graph Neural Network Cell Segmentation\nfor Spatial Transcriptomics",
             font_size=26, color=TEAL_ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
             "Technical Presentation  |  February 2026",
             font_size=18, color=RGBColor(180, 200, 220), alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: OUTLINE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Outline")
slide_number(sl, 2)

sections = [
    "1. Problem Formulation — Graph-based cell segmentation",
    "2. v0.1.0 Baseline — BCE loss and limitations",
    "3. v0.2.0 Architecture — Positional embeddings, Polars, cuSpatial",
    "4. Multi-Task Loss — Triplet, Metric, Alignment losses",
    "5. ME Gene Discovery — scRNA-seq biological priors",
    "6. Fragment Mode — GPU-accelerated connected components",
    "7. Adaptive Thresholding — Per-gene Li + Yen methods",
    "8. Delaunay Boundaries — De novo boundary extraction",
    "9. Checkpoint Workflows — Auto vocabulary management",
    "10. Performance & Scalability — Vectorization speedups",
]
add_multi_text(sl, Inches(1.2), Inches(1.2), Inches(10), Inches(6),
               sections, font_size=19, color=DARK_TEXT, line_spacing=1.5)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: SECTION - PROBLEM FORMULATION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "Problem Formulation", "Cell segmentation as link prediction on heterogeneous graphs")
slide_number(sl, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: CELL SEGMENTATION AS LINK PREDICTION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Cell Segmentation as Link Prediction")
slide_number(sl, 4)

add_multi_text(sl, Inches(0.6), Inches(1.1), Inches(7), Inches(1.2),
               ["Goal: Assign each RNA transcript t\u1d62 to its source cell c\u2c7c using spatial coordinates and gene identity."],
               font_size=18, bold_first=True)

add_equation_box(sl, Inches(0.6), Inches(2.2), Inches(7.5), Inches(0.7),
                 "\u0109(t\u1d62) = argmax_j  s\u1d62\u2c7c    s.t.  s\u1d62\u2c7c \u2265 \u03b8",
                 font_size=18)

add_text_box(sl, Inches(0.6), Inches(3.1), Inches(7.5), Inches(0.5),
             "where s\u1d62\u2c7c = \u27e8f(t\u1d62), f(c\u2c7c)\u27e9  (cosine similarity of GNN embeddings)",
             font_size=15, color=MED_GRAY)

add_multi_text(sl, Inches(0.6), Inches(3.8), Inches(6), Inches(3),
               ["Core insight: model segmentation as link prediction on a heterogeneous graph",
                "Graph: G = (V, E) where V = T \u222a C (transcripts \u222a cells)",
                "Node types: tx (gene tokens) and bd (boundary morphology)",
                "Similarity computed via L2-normalized GNN embeddings"],
               font_size=16, bullet=True)

# Info box on right
add_rect(sl, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.5), LIGHT_GRAY)
add_multi_text(sl, Inches(8.7), Inches(1.7), Inches(3.9), Inches(4),
               ["Boundary Features:",
                "Area: A(B\u1d62)",
                "Convexity: A(CH) / A(B\u1d62)",
                "Elongation: A(MBR) / A(Env)",
                "Circularity: A(B\u1d62) / r\u00b2_min",
                "",
                "Edge Construction:",
                "E_TT: KD-tree, d(p\u1d62,p\u2c7c) \u2264 r, |N(i)| \u2264 k",
                "E_TC: polygon containment",
                "E_BB: boundary spatial proximity"],
               font_size=14, color=TEAL_DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: v0.1.0 BASELINE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "v0.1.0 Baseline", "Binary cross-entropy link prediction")
slide_number(sl, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: BCE LOSS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "v0.1.0: BCE Loss for Link Prediction")
slide_number(sl, 6)

add_equation_box(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
                 "L_BCE = -\u03a3_(i,j) [ y\u1d62\u2c7c \u00b7 log \u03c3(s\u1d62\u2c7c) + (1 - y\u1d62\u2c7c) \u00b7 log(1 - \u03c3(s\u1d62\u2c7c)) ]",
                 font_size=18)

add_multi_text(sl, Inches(0.6), Inches(2.3), Inches(6), Inches(2.5),
               ["s\u1d62\u2c7c = \u27e8f(t\u1d62), f(c\u2c7c)\u27e9  (dot product of L2-normalized embeddings)",
                "y\u1d62\u2c7c = 1 if t\u1d62 belongs to c\u2c7c, else 0 (hard negatives)",
                "Hard negative sampling from nearby cells (1:5 ratio)",
                "SkipGAT architecture: GATv2Conv + HeteroConv"],
               font_size=16, bullet=True)

# Limitations box
add_rect(sl, Inches(0.6), Inches(4.6), Inches(12), Inches(2.5), RGBColor(255, 240, 240))
r = add_rect(sl, Inches(0.6), Inches(4.6), Inches(12), Inches(0.45), CORAL)
add_text_box(sl, Inches(0.8), Inches(4.65), Inches(5), Inches(0.4),
             "Limitations of v0.1.0", font_size=16, bold=True, color=WHITE)
add_multi_text(sl, Inches(0.8), Inches(5.2), Inches(5.5), Inches(2),
               ["Single loss \u2014 no embedding quality signal",
                "No biological priors (ME genes)",
                "Global similarity threshold \u03b8"],
               font_size=15, bullet=True, color=CORAL)
add_multi_text(sl, Inches(6.5), Inches(5.2), Inches(5.5), Inches(2),
               ["No fragment mode for unassigned transcripts",
                "Pandas/Dask: 2\u20133\u00d7 memory overhead",
                "buffer() \u2014 no polygon shrinking"],
               font_size=15, bullet=True, color=CORAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: SECTION - ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "v0.2.0 Architecture", "Positional embeddings, Polars, cuSpatial")
slide_number(sl, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: POSITIONAL EMBEDDINGS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Sinusoidal Positional Embeddings")
slide_number(sl, 8)

add_text_box(sl, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "Encode 2D spatial coordinates as frequency-domain features:",
             font_size=18, bold=True)

add_equation_box(sl, Inches(0.6), Inches(1.7), Inches(12), Inches(0.7),
                 "PE(p, 2k) = sin(p / 10000^(2k/d))     PE(p, 2k+1) = cos(p / 10000^(2k/d))",
                 font_size=17)

# Vectorization comparison
add_rect(sl, Inches(0.6), Inches(2.8), Inches(5.8), Inches(0.4), CORAL)
add_text_box(sl, Inches(0.8), Inches(2.83), Inches(5), Inches(0.35),
             "v0.1.0: Per-sample loop O(B \u00b7 N \u00b7 d)", font_size=14, bold=True, color=WHITE)
add_rect(sl, Inches(0.6), Inches(3.25), Inches(5.8), Inches(1.3), RGBColor(255, 240, 240))
add_text_box(sl, Inches(0.8), Inches(3.35), Inches(5.4), Inches(1.1),
             "for b in range(num_batches):\n"
             "    mask = (batch == b)\n"
             "    pos[mask] = normalize(pos[mask])",
             font_size=13, font_name="Courier New", color=DARK_TEXT)

add_rect(sl, Inches(6.8), Inches(2.8), Inches(5.8), Inches(0.4), TEAL_ACCENT)
add_text_box(sl, Inches(7.0), Inches(2.83), Inches(5), Inches(0.35),
             "v0.2.0: scatter ops O(N \u00b7 d)  \u2014  10\u2013100\u00d7 speedup", font_size=14, bold=True, color=WHITE)
add_rect(sl, Inches(6.8), Inches(3.25), Inches(5.8), Inches(1.3), RGBColor(230, 248, 245))
add_text_box(sl, Inches(7.0), Inches(3.35), Inches(5.4), Inches(1.1),
             "mins, _ = scatter_min(pos, batch)\n"
             "maxs, _ = scatter_max(pos, batch)\n"
             "pos_norm = (pos - mins[batch]) /\n"
             "           (maxs[batch] - mins[batch])",
             font_size=13, font_name="Courier New", color=DARK_TEXT)

# Pipeline diagram
add_text_box(sl, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
             "ISTEncoder Pipeline:", font_size=18, bold=True)

boxes = ["Gene\nEmbedding", "Positional\n2D Embed", "Linear\n\u2192 d_h", "SkipGAT\n\u00d7 (n+2)", "L2\nNormalize", "h \u2208 R^d_out"]
colors = [LIGHT_GRAY, RGBColor(235, 245, 255), RGBColor(235, 245, 255),
          RGBColor(255, 240, 230), RGBColor(255, 248, 220), RGBColor(230, 248, 245)]
for i, (txt, clr) in enumerate(zip(boxes, colors)):
    x = Inches(0.6 + i * 2.1)
    r = add_rect(sl, x, Inches(5.6), Inches(1.8), Inches(0.8), clr)
    r.line.color.rgb = TEAL_DARK
    r.line.width = Pt(1)
    add_text_box(sl, x + Inches(0.05), Inches(5.63), Inches(1.7), Inches(0.7),
                 txt, font_size=12, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        add_text_box(sl, x + Inches(1.85), Inches(5.8), Inches(0.3), Inches(0.4),
                     "\u2192", font_size=20, color=TEAL_DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 9: DATA PROCESSING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Data Processing: Polars & cuSpatial")
slide_number(sl, 9)

# Left column
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(0.45), TEAL_DARK)
add_text_box(sl, Inches(0.6), Inches(1.24), Inches(5), Inches(0.4),
             "Polars LazyFrames", font_size=16, bold=True, color=WHITE)
add_multi_text(sl, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3),
               ["Lazy evaluation \u2014 query optimization",
                "Streaming mode for >10M rows",
                "2\u20133\u00d7 less memory than Pandas",
                "Apache Arrow columnar backend",
                "Rust-based parallelism"],
               font_size=15, bullet=True)

# Right column
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(0.45), TEAL_DARK)
add_text_box(sl, Inches(7.0), Inches(1.24), Inches(5), Inches(0.4),
             "cuSpatial GPU Geometry", font_size=16, bold=True, color=WHITE)
add_multi_text(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3),
               ["Point-in-polygon (GPU-native)",
                "Quadtree spatial indexing",
                "Polygon area/centroid ops",
                "O(N/P) with P GPU cores"],
               font_size=15, bullet=True)

# Fallback table
add_rect(sl, Inches(0.4), Inches(5.0), Inches(12.4), Inches(0.4), DEEP_BLUE)
add_text_box(sl, Inches(0.6), Inches(5.03), Inches(5), Inches(0.35),
             "GPU / CPU Fallback Pattern", font_size=15, bold=True, color=WHITE)
fallback = [
    ("cuSpatial", "\u2192", "GeoPandas/Shapely"),
    ("CuPy sparse", "\u2192", "SciPy sparse"),
    ("cuML Louvain", "\u2192", "scanpy Leiden"),
]
for i, (gpu, arr, cpu) in enumerate(fallback):
    y = Inches(5.55 + i * 0.4)
    add_text_box(sl, Inches(2), y, Inches(3), Inches(0.35),
                 gpu, font_size=14, color=CORAL, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(sl, Inches(5.2), y, Inches(1), Inches(0.35),
                 arr, font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(6.5), y, Inches(3), Inches(0.35),
                 cpu, font_size=14, color=TEAL_DARK, bold=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: SECTION - MULTI-TASK LOSS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "Multi-Task Loss Formulation", "Triplet + Metric + Alignment with cosine scheduling")
slide_number(sl, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11: LOSS OVERVIEW
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "v0.2.0 Multi-Task Loss: Overview")
slide_number(sl, 11)

add_equation_box(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
                 "L_main = w_tx(t) \u00b7 L_tx  +  w_bd(t) \u00b7 L_bd  +  w_sg(t) \u00b7 L_sg",
                 font_size=20)

add_text_box(sl, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
             "Weight scheduling (smooth cosine annealing):", font_size=18, bold=True)

add_equation_box(sl, Inches(0.6), Inches(2.8), Inches(12), Inches(0.7),
                 "\u03b1(t) = \u00bd(1 + cos(\u03c0t/T))      w\u1d62(t) = w\u1d62_end + (w\u1d62_start \u2212 w\u1d62_end) \u00b7 \u03b1(t)",
                 font_size=17)

# Loss components
components = [
    ("L_tx", "Transcript Triplet Loss", "Cluster-aware embedding learning for transcripts", TEAL_ACCENT),
    ("L_bd", "Boundary Metric Loss", "Phenograph-supervised cell embedding quality", DEEP_BLUE),
    ("L_sg", "Segmentation Triplet Loss", "Link prediction between tx and bd nodes", CORAL),
    ("L_align", "Alignment Loss (optional)", "ME gene constraints from scRNA-seq reference", GOLD),
]
for i, (sym, name, desc, clr) in enumerate(components):
    y = Inches(4.0 + i * 0.8)
    r = add_rect(sl, Inches(0.6), y, Inches(1.2), Inches(0.6), clr)
    add_text_box(sl, Inches(0.7), y + Inches(0.1), Inches(1), Inches(0.4),
                 sym, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
                 font_name="Cambria Math")
    add_text_box(sl, Inches(2.0), y + Inches(0.02), Inches(4), Inches(0.3),
                 name, font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(sl, Inches(2.0), y + Inches(0.32), Inches(10), Inches(0.3),
                 desc, font_size=13, color=MED_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 12: TRIPLET LOSS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Triplet Loss with FastTripletSelector")
slide_number(sl, 12)

add_equation_box(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
                 "L_triplet = (1/|A|) \u03a3_(a,p,n) max(0,  \u2016h_a \u2212 h_p\u2016\u00b2 \u2212 \u2016h_a \u2212 h_n\u2016\u00b2 + m)",
                 font_size=18)

add_text_box(sl, Inches(0.6), Inches(2.2), Inches(12), Inches(0.4),
             "Margins: m_tx = 0.3, m_sg = 0.4.  Triplets sampled using phenograph cluster similarity.",
             font_size=15, color=MED_GRAY)

# FastTripletSelector
add_rect(sl, Inches(0.6), Inches(2.8), Inches(6), Inches(0.4), TEAL_DARK)
add_text_box(sl, Inches(0.8), Inches(2.83), Inches(5), Inches(0.35),
             "FastTripletSelector (O(1) per sample)", font_size=14, bold=True, color=WHITE)

steps = [
    "1. Build cluster similarity matrix S \u2208 R^(K\u00d7K) via phenograph",
    "2. Compute CDF:  pdf_k = S[a,k] / \u03a3_j S[a,j];  cdf = cumsum(pdf)",
    "3. Sample u ~ Unif(0,1)",
    "4. k* = searchsorted(cdf, u)  \u2014  O(log K) per sample",
    "5. Pick random node from cluster k*",
]
add_multi_text(sl, Inches(0.8), Inches(3.4), Inches(5.5), Inches(3.5),
               steps, font_size=14, color=DARK_TEXT)

# Right: similarity matrix concept
add_rect(sl, Inches(7.0), Inches(2.8), Inches(5.8), Inches(4.3), LIGHT_GRAY)
add_text_box(sl, Inches(7.2), Inches(2.9), Inches(5), Inches(0.4),
             "Cluster Similarity Matrix", font_size=16, bold=True, color=TEAL_DARK)
add_multi_text(sl, Inches(7.2), Inches(3.4), Inches(5.3), Inches(3.5),
               ["Phenograph/Louvain on boundary KNN graph",
                "S\u1d62\u2c7c = Jaccard similarity of clusters i, j",
                "Positive sampling: P(cluster k) \u221d S[a_cluster, k]",
                "Negative sampling: P(cluster k) \u221d \u2212S[a_cluster, k]",
                "",
                "Resolution \u03c1 = 2.0\u20133.0",
                "Inverse CDF: vectorized over all anchors"],
               font_size=14, color=DARK_TEXT, bullet=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 13: METRIC LOSS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Metric Loss: Cell Boundary Embedding Quality")
slide_number(sl, 13)

add_equation_box(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
                 "L_metric = E[ MSE(cos(h_a, h_p), 1\u2212d(a,p)) + MSE(cos(h_a, h_n), 1\u2212d(a,n)) ]",
                 font_size=18)

add_text_box(sl, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
             "where d(a,p) = 1 \u2212 S[c_a, c_p] is the cluster dissimilarity from phenograph.",
             font_size=16, color=MED_GRAY)

add_multi_text(sl, Inches(0.6), Inches(3.0), Inches(5.5), Inches(2),
               ["Intuition:",
                "Similar cells \u2192 high cosine similarity",
                "Dissimilar cells \u2192 low cosine similarity",
                "Supervised by phenograph clustering as soft labels"],
               font_size=16, bullet=True, bold_first=True)

add_multi_text(sl, Inches(7), Inches(3.0), Inches(5.5), Inches(2.5),
               ["Phenograph Pipeline:",
                "Compute cell-gene count vectors",
                "PCA reduction \u2192 R^d",
                "KNN graph (k = 30)",
                "Louvain community detection",
                "S: Jaccard between communities"],
               font_size=16, bullet=True, bold_first=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 14: ALIGNMENT LOSS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Alignment Loss: ME Gene Constraints")
slide_number(sl, 14)

add_equation_box(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
                 "L_align = \u03a3_(i,j)\u2208P (1 \u2212 s\u1d62\u2c7c)\u00b2  +  \u03a3_(i,j)\u2208N max(s\u1d62\u2c7c \u2212 m, 0)\u00b2",
                 font_size=18)

add_multi_text(sl, Inches(0.6), Inches(2.3), Inches(6), Inches(2),
               ["s\u1d62\u2c7c = \u27e8h\u1d62, h\u2c7c\u27e9 (cosine sim, L2-normalized)",
                "P: same-gene tx neighbors (attract)",
                "N: ME gene pairs (repel)",
                "Fixed margin m = 0.2",
                "|P| \u2264 3|N| (subsample positives)"],
               font_size=16, bullet=True)

# Bio motivation box
add_rect(sl, Inches(7), Inches(2.2), Inches(5.8), Inches(2.2), RGBColor(230, 248, 245))
add_rect(sl, Inches(7), Inches(2.2), Inches(5.8), Inches(0.4), TEAL_ACCENT)
add_text_box(sl, Inches(7.2), Inches(2.23), Inches(5), Inches(0.35),
             "Biological Motivation", font_size=14, bold=True, color=WHITE)
add_text_box(sl, Inches(7.2), Inches(2.7), Inches(5.4), Inches(1.6),
             "Mutually exclusive genes (markers of different cell types) "
             "should never co-localize in the same cell. Penalizing high "
             "similarity between ME neighbors discourages false cell merges "
             "\u2192 lower MECR.",
             font_size=14, color=TEAL_DARK)

# Vectorization
add_text_box(sl, Inches(0.6), Inches(4.8), Inches(12), Inches(0.4),
             "Vectorized ME Gene Matching (10\u201350\u00d7 speedup):", font_size=18, bold=True)

add_equation_box(sl, Inches(0.6), Inches(5.4), Inches(12), Inches(0.7),
                 "key(g_a, g_b) = min(g_a, g_b) \u00b7 G_max + max(g_a, g_b)    \u2192  torch.isin(edge_keys, me_keys)",
                 font_size=16)

add_text_box(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
             "Cantor-like pairing reduces O(|E| \u00b7 |ME|) loop to O(|E| + |ME|) hash lookup.",
             font_size=14, color=MED_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 15: LOSS COMBINATION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Loss Combination Modes")
slide_number(sl, 15)

# Interpolate
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(3.2), RGBColor(235, 245, 255))
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(0.45), TEAL_DARK)
add_text_box(sl, Inches(0.6), Inches(1.24), Inches(5), Inches(0.4),
             "Interpolate Mode (Default)", font_size=16, bold=True, color=WHITE)
add_equation_box(sl, Inches(0.6), Inches(1.85), Inches(5.6), Inches(0.55),
                 "L = (1 \u2212 w_align) \u00b7 L_main + w_align \u00b7 L_align", font_size=16)
add_multi_text(sl, Inches(0.6), Inches(2.6), Inches(5.6), Inches(1.5),
               ["Total scale \u2248 constant",
                "Smooth trade-off",
                "Recommended for most cases"],
               font_size=15, bullet=True)

# Additive
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(3.2), RGBColor(255, 240, 240))
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(0.45), CORAL)
add_text_box(sl, Inches(7.0), Inches(1.24), Inches(5), Inches(0.4),
             "Additive Mode", font_size=16, bold=True, color=WHITE)
add_equation_box(sl, Inches(7.0), Inches(1.85), Inches(5.6), Inches(0.55),
                 "L = L_main + w_align \u00b7 L_align", font_size=16)
add_multi_text(sl, Inches(7.0), Inches(2.6), Inches(5.6), Inches(1.5),
               ["Total scale increases over training",
                "May require LR adjustment",
                "For aggressive ME enforcement"],
               font_size=15, bullet=True)

# Full expansion
add_text_box(sl, Inches(0.4), Inches(4.8), Inches(12), Inches(0.4),
             "Full v0.2.0 loss expansion:", font_size=18, bold=True)
add_equation_box(sl, Inches(0.4), Inches(5.3), Inches(12.4), Inches(0.8),
                 "L_v2 = w_tx \u00b7 L_triplet^tx + w_bd \u00b7 L_metric^bd + w_sg \u00b7 L_triplet^sg + w_align \u00b7 L_align",
                 font_size=20)
add_text_box(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "  transcript        boundary        segmentation        ME constraint",
             font_size=13, color=MED_GRAY, font_name="Courier New")

# ══════════════════════════════════════════════════════════════
# SLIDE 16: SECTION - ME GENE DISCOVERY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "ME Gene Discovery", "Biological priors from scRNA-seq reference")
slide_number(sl, 16)

# ══════════════════════════════════════════════════════════════
# SLIDE 17: ME GENE DISCOVERY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "ME Gene Discovery from scRNA-seq Reference")
slide_number(sl, 17)

steps_left = [
    "Step 1: Marker Gene Identification",
    "\u2022 Rank genes by mean expression per cell type",
    "\u2022 Select top 10% (pos_percentile)",
    "\u2022 Filter: \u226530% of cells express gene",
    "",
    "Step 2: Mutual Exclusivity Filter",
    "\u2022 expr_in(g, c) > 0.25 (expressed in own type)",
    "\u2022 expr_out(g, c) < 0.03 (absent elsewhere)",
    "",
    "Step 3: Cross-Type Pairing",
    "\u2022 ME = {(g1, g2) : g1 \u2208 M_c1, g2 \u2208 M_c2, c1 \u2260 c2}",
]
add_multi_text(sl, Inches(0.6), Inches(1.1), Inches(6.5), Inches(5.5),
               steps_left, font_size=15)

# MECR box
add_rect(sl, Inches(7.5), Inches(1.2), Inches(5.3), Inches(2.5), LIGHT_GRAY)
add_rect(sl, Inches(7.5), Inches(1.2), Inches(5.3), Inches(0.4), TEAL_DARK)
add_text_box(sl, Inches(7.7), Inches(1.23), Inches(5), Inches(0.35),
             "MECR Validation Metric", font_size=14, bold=True, color=WHITE)
add_equation_box(sl, Inches(7.7), Inches(1.8), Inches(4.9), Inches(0.55),
                 "MECR(g1,g2) = P(g1 \u2227 g2) / P(g1 \u2228 g2)", font_size=16)
add_multi_text(sl, Inches(7.7), Inches(2.5), Inches(4.9), Inches(1),
               ["Lower is better", "Good: < 0.15", "High \u2192 cell merging artifacts"],
               font_size=14, bullet=True, color=TEAL_DARK)

# Pipeline arrow
add_rect(sl, Inches(7.5), Inches(4.0), Inches(5.3), Inches(3), RGBColor(255, 248, 220))
pipeline = ["scRNA-seq reference", "\u2193", "Marker genes per cell type",
            "\u2193", "ME exclusivity filter", "\u2193",
            "(g\u1d62, g\u2c7c) ME pairs", "\u2193", "L_align training signal"]
add_multi_text(sl, Inches(7.8), Inches(4.1), Inches(4.9), Inches(2.8),
               pipeline, font_size=14, color=DEEP_BLUE, line_spacing=1.15)

# ══════════════════════════════════════════════════════════════
# SLIDE 18: SECTION - FRAGMENT MODE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "Fragment Mode", "GPU-accelerated connected components for unassigned transcripts")
slide_number(sl, 18)

# ══════════════════════════════════════════════════════════════
# SLIDE 19: FRAGMENT MODE ALGORITHM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Fragment Mode: Recovering Unassigned Transcripts")
slide_number(sl, 19)

add_text_box(sl, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "Problem: Transcripts with max_j s\u1d62\u2c7c < \u03b8 remain unassigned after initial segmentation.",
             font_size=16, color=CORAL, bold=True)

steps = [
    "1. Identify unassigned transcripts:  T_u = {t\u1d62 : seg_cell_id(t\u1d62) = \u2205}",
    "2. Filter tx\u2013tx edges:  E_u = {(t\u1d62, t\u2c7c) \u2208 E_TT : t\u1d62, t\u2c7c \u2208 T_u}",
    "3. Threshold by similarity:  E'_u = {(t\u1d62, t\u2c7c) \u2208 E_u : s\u1d62\u2c7c \u2265 \u03b8_f}",
    "4. Connected components:  {C\u2081, C\u2082, ...} = CC(T_u, E'_u)",
    "5. Filter by size:  keep C_k if |C_k| \u2265 n_min (default: 5)",
    '6. Assign:  seg_cell_id(t\u1d62) = "fragment-k" for t\u1d62 \u2208 C_k',
]
add_multi_text(sl, Inches(0.6), Inches(1.8), Inches(12), Inches(3.5),
               steps, font_size=16, line_spacing=1.45)

# GPU vs CPU
add_rect(sl, Inches(0.4), Inches(5.2), Inches(6), Inches(2), RGBColor(255, 240, 230))
add_rect(sl, Inches(0.4), Inches(5.2), Inches(6), Inches(0.38), CORAL)
add_text_box(sl, Inches(0.6), Inches(5.22), Inches(5), Inches(0.35),
             "GPU: RAPIDS/CuPy", font_size=13, bold=True, color=WHITE)
add_text_box(sl, Inches(0.6), Inches(5.7), Inches(5.6), Inches(1.3),
             "CuPy CSR matrix \u2192 cupyx CC\n"
             "Parallel component labeling\n"
             "cp.bincount for size filtering",
             font_size=13, font_name="Courier New")

add_rect(sl, Inches(6.8), Inches(5.2), Inches(6), Inches(2), RGBColor(230, 248, 245))
add_rect(sl, Inches(6.8), Inches(5.2), Inches(6), Inches(0.38), TEAL_ACCENT)
add_text_box(sl, Inches(7.0), Inches(5.22), Inches(5), Inches(0.35),
             "CPU: SciPy Fallback", font_size=13, bold=True, color=WHITE)
add_text_box(sl, Inches(7.0), Inches(5.7), Inches(5.6), Inches(1.3),
             "scipy.sparse.csr_matrix\n"
             "scipy connected_components\n"
             "Complexity: O(|T_u| + |E'_u|)",
             font_size=13, font_name="Courier New")

# ══════════════════════════════════════════════════════════════
# SLIDE 20: ADAPTIVE THRESHOLDING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Per-Gene Adaptive Similarity Thresholding")
slide_number(sl, 20)

# Fixed
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(2.3), LIGHT_GRAY)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(0.4), DEEP_BLUE)
add_text_box(sl, Inches(0.6), Inches(1.23), Inches(5), Inches(0.35),
             "Fixed Threshold", font_size=15, bold=True, color=WHITE)
add_equation_box(sl, Inches(0.6), Inches(1.8), Inches(5.6), Inches(0.55),
                 "\u0109(t\u1d62) = argmax_j s\u1d62\u2c7c  if  max_j s\u1d62\u2c7c \u2265 \u03b8", font_size=15)
add_text_box(sl, Inches(0.6), Inches(2.5), Inches(5.6), Inches(0.7),
             "Single \u03b8 for all genes.\nSimple but may miss gene-specific variation.",
             font_size=14, color=MED_GRAY)

# Per-gene
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(2.3), RGBColor(255, 248, 220))
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(0.4), GOLD)
add_text_box(sl, Inches(7.0), Inches(1.23), Inches(5), Inches(0.35),
             "Per-Gene Auto-Threshold (v0.2.0)", font_size=15, bold=True, color=DARK_TEXT)
add_equation_box(sl, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.55),
                 "\u03b8_g = min(\u03b8_g^Li, \u03b8_g^Yen)   fallback: median(s_g)", font_size=15)
add_text_box(sl, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.7),
             "Li: minimum cross-entropy method\nYen: maximum correlation method\nRobust to multi-modal distributions.",
             font_size=14, color=MED_GRAY)

# Memory
add_text_box(sl, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
             "Memory optimization:  |S_g| > 10\u2077  \u2192  subsample to 10\u2077 before thresholding",
             font_size=16, color=TEAL_DARK, bold=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 21: SECTION - DELAUNAY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "Delaunay Boundary\nIdentification", "De novo cell boundary extraction from transcript clouds")
slide_number(sl, 21)

# ══════════════════════════════════════════════════════════════
# SLIDE 22: DELAUNAY ALGORITHM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Two-Phase Delaunay Boundary Extraction")
slide_number(sl, 22)

# Phase 1
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(2.8), RGBColor(235, 245, 255))
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6), Inches(0.4), DEEP_BLUE)
add_text_box(sl, Inches(0.6), Inches(1.23), Inches(5), Inches(0.35),
             "Phase 1: Remove Long Boundary Edges", font_size=15, bold=True, color=WHITE)
add_multi_text(sl, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2),
               ["Compute d_max = max nearest-neighbor distance",
                "Identify boundary edges (< 2 incident simplices)",
                "Iteratively remove edges with l > 2 \u00b7 d_max",
                "Propagate: deleted simplices \u2192 new boundary edges"],
               font_size=14, bullet=True)

# Phase 2
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(2.8), RGBColor(255, 240, 230))
add_rect(sl, Inches(6.8), Inches(1.2), Inches(6), Inches(0.4), CORAL)
add_text_box(sl, Inches(7.0), Inches(1.23), Inches(5), Inches(0.35),
             "Phase 2: Remove Extreme-Angle Edges", font_size=15, bold=True, color=WHITE)
add_multi_text(sl, Inches(7.0), Inches(1.8), Inches(5.6), Inches(2),
               ["For each boundary edge with opposite angle \u03b1:",
                "Remove if l > 1.5\u00b7d_max AND \u03b1 > 90\u00b0",
                "OR if \u03b1 > 180\u00b0 \u2212 180\u00b0/16 \u2248 168.75\u00b0",
                "Vectorized angle computation for all triangles"],
               font_size=14, bullet=True)

# Cycle detection
add_text_box(sl, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
             "Cycle Detection \u2192 Polygon Extraction:", font_size=18, bold=True)
add_multi_text(sl, Inches(0.6), Inches(4.9), Inches(12), Inches(1.5),
               ["DFS on remaining boundary graph to find connected components",
                "Extract cycles \u2192 Shapely Polygon / MultiPolygon",
                "Complexity: O(n log n) dominated by Delaunay triangulation"],
               font_size=15, bullet=True)

# Angle formula
add_equation_box(sl, Inches(0.4), Inches(6.3), Inches(12.4), Inches(0.7),
                 "\u03b1_k = arccos( (u \u00b7 v) / (\u2016u\u2016 \u00b7 \u2016v\u2016) )    \u2014  vectorized over all triangles simultaneously",
                 font_size=15)

# ══════════════════════════════════════════════════════════════
# SLIDE 23: CHECKPOINT WORKFLOWS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Simplified Checkpoint Workflows")
slide_number(sl, 23)

# Table
modes = [
    ("train", "Fresh training", "Embeddings from data module vocab", TEAL_ACCENT),
    ("finetune", "Load weights, reset optimizer", "Auto-remap: shared genes keep weights,\nnew genes Xavier init", DEEP_BLUE),
    ("predict", "Skip training", "Auto-filter: data restricted to\ncheckpoint gene vocabulary", CORAL),
]

# Header
y = Inches(1.3)
for x, w, label in [(Inches(0.6), Inches(2.2), "Mode"),
                     (Inches(2.8), Inches(4), "Description"),
                     (Inches(6.8), Inches(6), "Vocabulary Handling")]:
    r = add_rect(sl, x, y, w, Inches(0.45), DEEP_BLUE)
    add_text_box(sl, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.35),
                 label, font_size=14, bold=True, color=WHITE)

for i, (mode, desc, vocab, clr) in enumerate(modes):
    y = Inches(1.85 + i * 1.0)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(sl, Inches(0.6), y, Inches(2.2), Inches(0.85), clr)
    add_text_box(sl, Inches(0.7), y + Inches(0.2), Inches(2), Inches(0.5),
                 mode, font_size=16, bold=True, color=WHITE, font_name="Courier New")
    add_rect(sl, Inches(2.8), y, Inches(4), Inches(0.85), bg)
    add_text_box(sl, Inches(2.9), y + Inches(0.15), Inches(3.8), Inches(0.6),
                 desc, font_size=14, color=DARK_TEXT)
    add_rect(sl, Inches(6.8), y, Inches(6), Inches(0.85), bg)
    add_text_box(sl, Inches(6.9), y + Inches(0.08), Inches(5.8), Inches(0.7),
                 vocab, font_size=13, color=TEAL_DARK)

# Vocab remapping
add_text_box(sl, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             "Vocabulary Remapping (Finetune Mode):", font_size=17, bold=True)
add_equation_box(sl, Inches(0.6), Inches(5.7), Inches(12), Inches(0.7),
                 "V_shared = V_ckpt \u2229 V_data    V_new = V_data \\ V_ckpt    (shared: copy weights, new: Xavier init)",
                 font_size=15)
add_text_box(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             "Gene embeddings always exported as gene_embeddings.parquet for downstream reuse.",
             font_size=14, color=MED_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 24: SECTION - PERFORMANCE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(sl, "Performance & Scalability", "Vectorization speedups and GPU/CPU dual-path")
slide_number(sl, 24)

# ══════════════════════════════════════════════════════════════
# SLIDE 25: VECTORIZATION SPEEDUPS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Vectorization Speedups")
slide_number(sl, 25)

rows = [
    ("Batch normalization", "Per-sample loop O(BNd)", "scatter_min/max O(Nd)", "10\u2013100\u00d7"),
    ("ME gene matching", "Per-edge loop O(|E|\u00b7|ME|)", "Cantor hash + isin O(|E|+|ME|)", "10\u201350\u00d7"),
    ("Triplet sampling", "Random rejection O(NK)", "FastTripletSelector O(N)", "~5\u00d7"),
    ("Data loading", "Pandas + Dask", "Polars LazyFrames (streaming)", "2\u20133\u00d7"),
]

# Header
y0 = Inches(1.3)
headers = ["Operation", "v0.1.0", "v0.2.0", "Speedup"]
widths = [Inches(2.5), Inches(3.5), Inches(3.5), Inches(1.8)]
x = Inches(0.6)
for w, label in zip(widths, headers):
    r = add_rect(sl, x, y0, w, Inches(0.45), DEEP_BLUE)
    add_text_box(sl, x + Inches(0.1), y0 + Inches(0.05), w - Inches(0.2), Inches(0.35),
                 label, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    x += w

for i, (op, old, new, speedup) in enumerate(rows):
    y = Inches(1.85 + i * 0.75)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    x = Inches(0.6)
    vals = [op, old, new, speedup]
    for j, (w, val) in enumerate(zip(widths, vals)):
        add_rect(sl, x, y, w, Inches(0.65), bg)
        clr = CORAL if j == 3 else DARK_TEXT
        bld = j == 3
        add_text_box(sl, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2), Inches(0.45),
                     val, font_size=13, color=clr, bold=bld, alignment=PP_ALIGN.CENTER)
        x += w

# Additional features
add_text_box(sl, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             "Additional scalability features:", font_size=17, bold=True)
add_multi_text(sl, Inches(0.6), Inches(5.7), Inches(6), Inches(1.5),
               ["Mixed precision: 16-mixed, bf16-mixed",
                "Auto 3D detection from z-coordinates",
                "Chunked GPU similarity (256MB target)"],
               font_size=15, bullet=True)
add_multi_text(sl, Inches(7), Inches(5.7), Inches(6), Inches(1.5),
               ["LR schedulers: cosine, OneCycle",
                "Early stopping with configurable patience",
                "Tile-based: 50K\u2013100K tx/tile"],
               font_size=15, bullet=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 26: COMPARISON TABLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "v0.1.0 vs. v0.2.0: Complete Comparison")
slide_number(sl, 26)

comparison = [
    ("Data processing", "Pandas + Dask", "Polars (2\u20133\u00d7 faster)"),
    ("GPU geometry", "RAPIDS cuML", "cuSpatial (native polygon)"),
    ("Loss function", "BCE only", "BCE + Triplet + Metric + Alignment"),
    ("Triplet sampling", "Random", "FastTripletSelector (inverse CDF)"),
    ("Biological priors", "None", "ME gene alignment loss"),
    ("Gene embeddings", "One-hot / external", "Learned + PCA + exported"),
    ("Positional encoding", "Per-sample norm", "scatter_min/max (10\u2013100\u00d7)"),
    ("Polygon scaling", "buffer() only", "scale() (expand + shrink)"),
    ("Similarity threshold", "Global \u03b8", "Per-gene Li + Yen auto"),
    ("Unassigned transcripts", "Discarded", "Fragment mode (GPU CC)"),
    ("Boundary construction", "External only", "Delaunay (de novo)"),
    ("Checkpoint system", "Manual vocab", "Auto (train/finetune/predict)"),
    ("Export formats", "Parquet", "Zarr, SpatialData, AnnData, Explorer"),
    ("Training precision", "FP32 only", "FP16/BF16 mixed precision"),
]

# Header
y0 = Inches(1.1)
hdrs = [("Feature", Inches(2.8)), ("v0.1.0", Inches(3.5)), ("v0.2.0", Inches(5.8))]
x = Inches(0.4)
for label, w in hdrs:
    add_rect(sl, x, y0, w, Inches(0.35), DEEP_BLUE)
    add_text_box(sl, x + Inches(0.08), y0 + Inches(0.02), w - Inches(0.16), Inches(0.3),
                 label, font_size=12, bold=True, color=WHITE)
    x += w

for i, (feat, old, new) in enumerate(comparison):
    y = Inches(1.5 + i * 0.39)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    x = Inches(0.4)
    for val, w in [(feat, Inches(2.8)), (old, Inches(3.5)), (new, Inches(5.8))]:
        add_rect(sl, x, y, w, Inches(0.36), bg)
        clr = TEAL_DARK if val == new and val != old else DARK_TEXT
        bld = val == new and val != old
        add_text_box(sl, x + Inches(0.08), y + Inches(0.03), w - Inches(0.16), Inches(0.3),
                     val, font_size=11, color=clr, bold=bld)
        x += w

# ══════════════════════════════════════════════════════════════
# SLIDE 27: KEY EQUATIONS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, "Key Equations Summary")
slide_number(sl, 27)

equations = [
    ("(1) Assignment", "\u0109(t\u1d62) = argmax_j \u27e8f(t\u1d62), f(c\u2c7c)\u27e9   s.t.  s\u1d62\u2c7c \u2265 \u03b8_g"),
    ("(2) BCE", "L_BCE = \u2212\u03a3 [y\u1d62\u2c7c log\u03c3(s\u1d62\u2c7c) + (1\u2212y\u1d62\u2c7c)log(1\u2212\u03c3(s\u1d62\u2c7c))]"),
    ("(3) Triplet", "L_tri = (1/|A|) \u03a3 max(0, \u2016h_a\u2212h_p\u2016\u00b2 \u2212 \u2016h_a\u2212h_n\u2016\u00b2 + m)"),
    ("(4) Metric", "L_met = E[MSE(cos(h_a,h_p), 1\u2212d_ap) + MSE(cos(h_a,h_n), 1\u2212d_an)]"),
    ("(5) Alignment", "L_align = \u03a3_P (1\u2212s\u1d62\u2c7c)\u00b2 + \u03a3_N max(s\u1d62\u2c7c\u2212m, 0)\u00b2"),
    ("(6) Combined", "L = (1\u2212w_a)[w_t\u00b7L_tri^t + w_b\u00b7L_met^b + w_s\u00b7L_tri^s] + w_a\u00b7L_align"),
    ("(7) Scheduling", "w\u1d62(t) = w\u1d62_end + (w\u1d62_start \u2212 w\u1d62_end) \u00b7 \u00bd(1 + cos(\u03c0t/T))"),
]

for i, (label, eq) in enumerate(equations):
    y = Inches(1.1 + i * 0.85)
    clr = [TEAL_ACCENT, MED_GRAY, TEAL_DARK, DEEP_BLUE, CORAL, DEEP_BLUE, GOLD][i]
    add_rect(sl, Inches(0.4), y, Inches(2), Inches(0.6), clr)
    add_text_box(sl, Inches(0.5), y + Inches(0.12), Inches(1.8), Inches(0.4),
                 label, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_equation_box(sl, Inches(2.5), y, Inches(10.2), Inches(0.6), eq, font_size=14)

# ══════════════════════════════════════════════════════════════
# SLIDE 28: THANK YOU
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, DEEP_BLUE)
add_rect(sl, Inches(0), Inches(5.8), W, Inches(0.06), TEAL_ACCENT)

add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Thank You", font_size=52, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "Questions?", font_size=32, color=TEAL_ACCENT,
             alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(2), Inches(4.8), Inches(4.5), Inches(1.5),
             "Code:  segger v0.2.0\n"
             "Install:  pip install segger",
             font_size=18, color=RGBColor(180, 200, 220))
add_text_box(sl, Inches(7), Inches(4.8), Inches(4.5), Inches(1.5),
             "CLI:\n"
             "segger segment -i data/ -o out/\n"
             "segger export -s seg.pq -o exp/",
             font_size=16, color=RGBColor(180, 200, 220), font_name="Courier New")

add_text_box(sl, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Segger: scalable graph neural network cell segmentation (2025)",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────
out_path = "/Users/b450-admin/Desktop/Projects/Active/segger/segger-0.2.0/presentation/segger_v2_technical.pptx"
prs.save(out_path)
print(f"Saved PPTX: {out_path}")
print(f"Total slides: {len(prs.slides)}")
